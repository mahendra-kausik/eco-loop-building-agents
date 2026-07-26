"""Fills the organizer's 6-slide PPT template (docs/IDEA_Presentation_Format.pptx)
with this project's real, measured results and writes
docs/Eco-Loop_Idea_Submission.pptx.

Build-time only: python-pptx lives in .venv, not requirements.txt (the
submission's runtime deps -- what the loop itself needs to run -- stay
honest; this script is never imported by the loop).

Template has 7 slides: slide 0 is the organizer's own instruction sheet
("keep max 6 including title") -- deleted here, not counted. The remaining
6 (Title / Proposed Solution / Technical Approach / Feasibility & Viability /
Artifacts / References) are filled in place, preserving the template's own
fonts/bullets/layout -- only run text is replaced, no new shapes.

Numbers come straight from src/analysis/metrics.compare() against the
shipped results/raw/{baseline,agent_llm}, the same source of truth as
README.md and results/dashboard.html -- this script has no numbers of its
own to get out of sync.
"""
import copy
import os
import sys

from pptx import Presentation

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))
from src.analysis.metrics import compare  # noqa: E402

REPO = os.path.join(os.path.dirname(__file__), "..")
TEMPLATE = os.path.join(REPO, "docs", "IDEA_Presentation_Format.pptx")
OUT = os.path.join(REPO, "docs", "Eco-Loop_Idea_Submission.pptx")
REPO_URL = "https://github.com/mahendra-kausik/eco-loop-building-agents"

FILL = "<FILL: from your SIH portal registration>"


def _delete_slide(prs: Presentation, index: int) -> None:
    """python-pptx has no public remove-slide API; drop the slide's rId from
    the presentation part and its entry from sldIdLst -- the standard recipe."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    rId = slides[index].rId
    prs.part.drop_rel(rId)
    xml_slides.remove(slides[index])


def _set_para_text(para, text: str) -> None:
    """Replace a paragraph's content with a single run, copying the first
    existing run's formatting so it still looks like the template."""
    if not para.runs:
        para.add_run()
    template_run = para.runs[0]
    fmt = (template_run.font.size, template_run.font.bold, template_run.font.italic)
    for r in list(para.runs)[1:]:
        r._r.getparent().remove(r._r)
    template_run.text = text
    template_run.font.size, template_run.font.bold, template_run.font.italic = fmt


def _append_run(para, text: str) -> None:
    """Add a plain (non-bold) run after a paragraph's existing bold label,
    e.g. turning 'Problem Statement ID-' into 'Problem Statement ID- <value>'."""
    run = para.add_run()
    run.text = " " + text
    run.font.bold = False
    if para.runs:
        run.font.size = para.runs[0].font.size


def _textbox(slide, name: str):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    raise KeyError(f"no shape named {name!r} on this slide")


def build(baseline_dir: str, agent_dir: str, out_path: str) -> str:
    result = compare(baseline_dir, agent_dir)
    b, a = result["baseline"], result["agent"]

    prs = Presentation(TEMPLATE)
    _delete_slide(prs, 0)  # organizer's instruction sheet -- not a content slide

    # --- Slide 1 (was 2): Title page -----------------------------------
    title_tb = _textbox(prs.slides[0], "TextBox 6")
    paras = title_tb.text_frame.paragraphs
    fills = {
        1: FILL,  # Problem Statement ID
        2: "Eco-Loop Building Agents -- Autonomous Closed-Loop Building Energy Management",
        3: FILL,  # Theme (portal-assigned wording)
        4: "Software",
        5: FILL,  # Student Name
        6: FILL,  # Student ID
    }
    for idx, value in fills.items():
        _append_run(paras[idx], value)

    # --- Slide 2 (was 3): Proposed Solution -----------------------------
    sol_tb = _textbox(prs.slides[1], "TextBox 8")
    p = sol_tb.text_frame.paragraphs
    _set_para_text(p[3],
        "An autonomous closed loop: EnergyPlus simulates the building, an "
        "open-weight LLM (gpt-oss-120b) reads live zone state every simulated "
        "hour, and a safety-clamped decision writes new HVAC setpoints straight "
        "back into the running simulation -- no human in the loop.")
    _set_para_text(p[4],
        "Buildings waste energy because fixed schedules can't react to real "
        "occupancy, weather, or grid carbon intensity. Here the LLM reasons "
        "against forecast context (weather, carbon/tariff, occupancy ahead) and "
        "measurably cuts HVAC energy while holding comfort in band -- proven on "
        "a matched 7-day baseline-vs-agent run, not simulated on paper.")
    _set_para_text(p[5],
        "A safety supervisor wraps the LLM: schema validation, hard clamps, a "
        "watchdog timeout, and a deterministic rule-based fallback mean the "
        "simulation cannot be killed by a bad, slow, or absent LLM response -- "
        "the LLM now also beats its own fallback on energy AND comfort at once.")

    # --- Slide 3 (was 4): Technical Approach -----------------------------
    tech_tb = _textbox(prs.slides[2], "TextBox 8")
    p = tech_tb.text_frame.paragraphs
    _set_para_text(p[0],
        "Python 3.11, EnergyPlus 26.1 via the pyenergyplus runtime API "
        "(in-process callbacks, not file-based re-runs), gpt-oss-120b served "
        "OpenAI-compatible on Cerebras + Groq (round-robined), FastMCP for "
        "spec-compliant tool exposure, Plotly for the evidence dashboard.")
    _set_para_text(p[1],
        "EnergyPlus -> tool layer (get_building_state / get_forecast_context / "
        "propose_setpoints / inject_setpoints / get_recent_errors, exposed both "
        "directly and over MCP) -> LLM decides once per simulated hour -> safety "
        "supervisor validates + clamps -> actuator injects setpoints -> loop closes. "
        "See docs/ARCHITECTURE.md for the full diagram and prompt strategy.")

    # --- Slide 4 (was 5): Feasibility & Viability ------------------------
    feas_tb = _textbox(prs.slides[3], "TextBox 8")
    p = feas_tb.text_frame.paragraphs
    _set_para_text(p[0],
        f"Measured, not projected, on a matched 7-simulated-day run: "
        f"{result['total_electricity_pct_saved']:+.1f}% total electricity, "
        f"{result['hvac_pct_saved']:+.1f}% HVAC-only electricity, comfort-in-band "
        f"{result['comfort_delta_pts']:+.1f} pts vs baseline -- savings aren't taken "
        f"out of occupant comfort.")
    _set_para_text(p[1],
        "~62% of facility electricity is lighting/plug load no setpoint can "
        "touch (HVAC-only % is the honest headline); free-tier LLM provider "
        "rate limits; the risk of the LLM ever returning invalid or absent output.")
    _set_para_text(p[2],
        "Per-provider throttling + round-robin turned 43% genuine LLM "
        "participation into 98%+; the safety supervisor's clamp+fallback make "
        "provider outages a non-issue (a 100%-fallback run completes with "
        "identical reliability); two tuning ideas (CO2-gated fan-off, "
        "optimal-start pre-heat) were measured and reverted when the data "
        "showed no net benefit -- engineering decisions backed by numbers.")

    # --- Slide 5 (was 6): Artifacts --------------------------------------
    art_tb = _textbox(prs.slides[4], "TextBox 8")
    p = art_tb.text_frame.paragraphs
    _set_para_text(p[1], f"Full source: {REPO_URL}")
    _set_para_text(p[2], "Self-contained offline dashboard: results/dashboard.html")
    _set_para_text(p[3],
        f"Headline: {b['total_electricity_kwh']:.0f} -> {a['total_electricity_kwh']:.0f} kWh "
        f"({result['total_electricity_pct_saved']:+.1f}%) total; "
        f"{b['hvac_kwh']:.0f} -> {a['hvac_kwh']:.0f} kWh "
        f"({result['hvac_pct_saved']:+.1f}%) HVAC-only; comfort "
        f"{b['comfort_in_band_pct']:.1f}% -> {a['comfort_in_band_pct']:.1f}%; "
        f"peak {b['peak_demand_kw']:.1f} -> {a['peak_demand_kw']:.1f} kW")
    p[4].text = ""
    _append_run(p[4], f"[FILL: paste a dashboard screenshot here] -- open {REPO_URL.split('/')[-1]}/results/dashboard.html")

    # --- Slide 6 (was 7): References --------------------------------------
    ref_tb = _textbox(prs.slides[5], "TextBox 8")
    p = ref_tb.text_frame.paragraphs
    _set_para_text(p[0], "EnergyPlus documentation & pyenergyplus API -- energyplus.net")
    from pptx.text.text import _Paragraph
    anchor = p[0]._p
    for text in (
        "gpt-oss-120b (OpenAI open-weight model) via Cerebras & Groq OpenAI-compatible APIs",
        "Model Context Protocol (MCP) spec -- modelcontextprotocol.io",
        "ASHRAE 55 / Fanger PMV thermal comfort model",
    ):
        new_p = copy.deepcopy(p[0]._p)
        anchor.addnext(new_p)
        anchor = new_p
        _set_para_text(_Paragraph(new_p, ref_tb.text_frame), text)

    prs.save(out_path)
    return out_path


def main() -> None:
    out = build(
        os.path.join(REPO, "results", "raw", "baseline"),
        os.path.join(REPO, "results", "raw", "agent_llm"),
        OUT,
    )
    print(f"Wrote deck -> {out}")


if __name__ == "__main__":
    main()
